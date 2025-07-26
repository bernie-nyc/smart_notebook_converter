"""
Module to convert SMART Notebook (.notebook) files into PowerPoint (.pptx).

SMART Notebook files are essentially ZIP archives containing page
descriptions in vector (SVG) and bitmap (PNG/JPEG) formats along with
metadata.  This script extracts each supported page asset—SVG pages
are converted to PNG; PNG and JPEG pages are used directly—and then
composes a PowerPoint presentation by placing each image on its own
slide.  Interactive elements such as audio, video, or embedded
activities are intentionally ignored; the goal is to produce a
static slide deck.

Dependencies:

* python-pptx – used to create the output presentation.  This library
  is available in the current environment.
* A tool to convert SVG to PNG.  The converter attempts to use
  `cairosvg` (if installed) and falls back to ImageMagick’s `magick`
  command.  On systems without either of these, you must install
  CairoSVG or ensure that ImageMagick is properly configured with
  delegates for SVG.

Usage:

    python notebook_to_ppt.py --input <path_to_notebook_or_directory> --output-dir <output_directory>

If the input is a directory, the script will recursively scan for
.notebook files.  For each file it finds, it will produce a .pptx
file with the same base name in the specified output directory.

Limitations:

* SVG, PNG and JPEG pages are processed.  Some Notebook files may
  store page content in other proprietary formats (for example, as
  custom XML).  Those unsupported pages will be skipped.
* Interactive elements, audio, and video are not transferred.
* Conversion quality depends on the SVG converter; complex pages may
  render differently from the original Notebook.
"""

from __future__ import annotations

import argparse
import logging
import os
import subprocess
import sys
import tempfile
import zipfile
import shutil

from pathlib import Path
from typing import Iterable, List

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as exc:
    raise SystemExit("python-pptx is required to run this script") from exc

_log = logging.getLogger(__name__)

def convert_svg_to_png(svg_path: Path, png_path: Path) -> None:
    """Convert an SVG file to a PNG file.

    Attempts to use cairosvg first, then falls back to ImageMagick’s
    ``magick`` CLI.  Raises RuntimeError on failure.
    """
    try:
        import cairosvg  # type: ignore
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path))
        return
    except ImportError:
        pass
    except Exception as exc:
        _log.warning("cairosvg failed: %s", exc)

    magick_path = shutil.which("magick")
    if magick_path:
        try:
            subprocess.run(
                [magick_path, str(svg_path), str(png_path)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            return
        except subprocess.CalledProcessError as exc:
            _log.warning("ImageMagick failed: %s", exc.stderr.decode())

    raise RuntimeError(
        "Unable to convert SVG to PNG. Install cairosvg or ensure ImageMagick is configured."
    )

def extract_page_files(zf: zipfile.ZipFile) -> List[str]:
    """Return a list of page asset file names in the notebook archive."""
    candidates: List[str] = []
    for name in zf.namelist():
        base = os.path.basename(name).lower()
        if base.startswith('page') and (
            base.endswith('.svg')
            or base.endswith('.png')
            or base.endswith('.jpg')
            or base.endswith('.jpeg')
        ):
            candidates.append(name)
    def page_key(name: str) -> int:
        digits = ''.join(ch for ch in os.path.basename(name) if ch.isdigit())
        return int(digits) if digits else 0
    return sorted(candidates, key=page_key)

def process_notebook(notebook_path: Path, output_dir: Path) -> Path:
    """Convert a single .notebook file into a .pptx."""
    if not notebook_path.exists():
        raise FileNotFoundError(notebook_path)
    output_dir.mkdir(parents=True, exist_ok=True)
    base_name = notebook_path.stem
    pptx_path = output_dir / f"{base_name}.pptx"
    _log.info("Processing %s", notebook_path)
    with zipfile.ZipFile(notebook_path) as zf:
        page_files = extract_page_files(zf)
        if not page_files:
            _log.warning("No supported page assets found in %s; skipping", notebook_path)
            return pptx_path
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        with tempfile.TemporaryDirectory() as tmpdir_name:
            tmpdir = Path(tmpdir_name)
            for idx, page_name in enumerate(page_files, start=1):
                suffix = Path(page_name).suffix.lower()
                extracted_path = tmpdir / f"page_{idx}{suffix}"
                with extracted_path.open('wb') as f_out:
                    f_out.write(zf.read(page_name))
                if suffix == '.svg':
                    png_path = tmpdir / f"page_{idx}.png"
                    try:
                        convert_svg_to_png(extracted_path, png_path)
                    except Exception as exc:
                        _log.error("Failed to convert %s: %s", page_name, exc)
                        continue
                    image_path = png_path
                else:
                    image_path = extracted_path
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height
                )
            prs.save(pptx_path)
    _log.info("Saved PowerPoint to %s", pptx_path)
    return pptx_path

def iter_notebook_files(input_path: Path) -> Iterable[Path]:
    """Yield .notebook files from the given path."""
    if input_path.is_file() and input_path.suffix.lower() == '.notebook':
        yield input_path
    elif input_path.is_dir():
        for path in input_path.rglob('*.notebook'):
            yield path
    else:
        _log.warning("%s is neither a .notebook file nor a directory", input_path)

def main(argv: List[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        '--input', '-i', required=True, type=Path,
        help='Path to a .notebook file or directory containing .notebook files.'
    )
    parser.add_argument(
        '--output-dir', '-o', required=True, type=Path,
        help='Directory where converted .pptx files will be written.'
    )
    parser.add_argument(
        '--verbose', '-v', action='store_true', help='Enable verbose logging.'
    )
    args = parser.parse_args(argv)
    logging.basicConfig(level=logging.DEBUG if args.verbose else logging.INFO)
    any_processed = False
    for notebook_file in iter_notebook_files(args.input):
        any_processed = True
        try:
            process_notebook(notebook_file, args.output_dir)
        except Exception as exc:
            _log.error("Error processing %s: %s", notebook_file, exc)
    if not any_processed:
        _log.error("No .notebook files found in %s", args.input)
        return 1
    return 0

if __name__ == '__main__':
    raise SystemExit(main())
