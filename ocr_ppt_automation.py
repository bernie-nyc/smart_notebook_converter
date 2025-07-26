import os
import shutil
import comtypes.client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import pytesseract
from pytesseract import Output
from cairosvg import svg2png

# STEP 1: Convert each slide of a PowerPoint file into PNG images.
def ppt_to_images(input_ppt, output_dir):
    # Start PowerPoint through Windows COM automation
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1  # Make PowerPoint visible (required by COM API)

    # Open the presentation file
    presentation = powerpoint.Presentations.Open(input_ppt)

    # Export slides as PNG images into the output directory
    presentation.Export(output_dir, "PNG")

    # Close presentation and PowerPoint application
    presentation.Close()
    powerpoint.Quit()

# STEP 2: Convert SVG files to PNG if needed (Tesseract cannot process SVG).
def convert_svg_images(image_dir):
    for file in os.listdir(image_dir):
        if file.lower().endswith('.svg'):
            svg_path = os.path.join(image_dir, file)
            png_path = os.path.join(image_dir, file[:-4] + '.png')
            # Convert SVG to PNG using CairoSVG
            svg2png(url=svg_path, write_to=png_path)
            os.remove(svg_path)  # Remove original SVG after conversion

# STEP 3: Perform OCR and capture each text box with its size, position, and color.
def ocr_images_with_layout(image_path):
    # Set path to Tesseract executable
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

    # Open the image file
    img = Image.open(image_path)

    # Run OCR with layout data returned as a dictionary
    data = pytesseract.image_to_data(img, output_type=Output.DICT)

    elements = []  # Store structured text elements for slide reconstruction

    # Loop through all OCR-detected text boxes
    n_boxes = len(data['level'])
    for i in range(n_boxes):
        # Filter out low-confidence or empty text regions
        if int(data['conf'][i]) > 60 and data['text'][i].strip():
            x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]

            # Crop the image region for this text box to get average color
            cropped_img = img.crop((x, y, x + w, y + h))
            avg_color = cropped_img.convert('RGB').resize((1, 1)).getpixel((0, 0))  # Resize to 1x1 for color avg

            # Store bounding box and text data
            elements.append({
                'text': data['text'][i],  # The actual recognized text
                'x': x,                  # Horizontal position in pixels
                'y': y,                  # Vertical position in pixels
                'w': w,                  # Width of bounding box
                'h': h,                  # Height of bounding box
                'color': avg_color       # Average color (R,G,B) tuple
            })
    return elements

# STEP 4: Create one slide per image and place text boxes based on OCR positions.
def create_layout_slide(prs, elements):
    # Add a completely blank slide (no title, no content boxes)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Loop through each OCR-extracted text box and place it on the slide
    for elem in elements:
        # Convert pixel coordinates to inches (1 inch = 96 pixels)
        textbox = slide.shapes.add_textbox(
            Inches(elem['x'] / 96),  # Left
            Inches(elem['y'] / 96),  # Top
            Inches(elem['w'] / 96),  # Width
            Inches(elem['h'] / 96)   # Height
        )

        # Add the text into the textbox
        frame = textbox.text_frame
        frame.clear()  # Remove the default paragraph PowerPoint adds
        p = frame.paragraphs[0]  # Create a fresh paragraph for the text
        p.text = elem['text']

        # Estimate font size based on height of bounding box (scaled empirically)
        p.font.size = Pt(max(8, elem['h'] * 0.75))

        # Assign text color based on average color in that region
        p.font.color.rgb = RGBColor(*elem['color'])

# STEP 5: Full OCR pipeline for a single .pptx file
# Converts to images, OCRs layout, rebuilds into editable text slide

def process_ppt_file(input_ppt):
    base_name = os.path.splitext(input_ppt)[0]  # Remove .pptx extension
    temp_dir = base_name + "_temp_images"       # Temp directory for exported slide images
    output_ppt = base_name + "_ocr_output.pptx" # Final output file

    # Clean up old temp data if exists
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
    os.makedirs(temp_dir, exist_ok=True)

    # Convert slides to PNGs, handle SVGs if needed
    ppt_to_images(input_ppt, temp_dir)
    convert_svg_images(temp_dir)

    # Create a new blank PowerPoint
    prs = Presentation()

    # OCR each slide image and create one slide per result
    image_files = sorted([file for file in os.listdir(temp_dir) if file.lower().endswith('.png')])
    for img_file in image_files:
        img_path = os.path.join(temp_dir, img_file)
        elements = ocr_images_with_layout(img_path)  # OCR + layout info
        create_layout_slide(prs, elements)           # Add slide to PPT

    # Save final presentation
    prs.save(output_ppt)

    # Clean up temporary image files
    shutil.rmtree(temp_dir)

# STEP 6: Process all PowerPoint files recursively in the given directory tree
def process_all_ppts(root_dir):
    for subdir, _, files in os.walk(root_dir):
        for file in files:
            if file.lower().endswith('.pptx'):
                input_ppt = os.path.join(subdir, file)
                process_ppt_file(input_ppt)

# Entry point: start from the current working directory
if __name__ == "__main__":
    current_dir = os.getcwd()
    process_all_ppts(current_dir)
